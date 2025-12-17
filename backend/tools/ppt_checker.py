from langchain_core.tools import BaseTool
from pptx import Presentation
import os

class PPTCheckTool(BaseTool):
    name: str = "ppt_check_tool"
    description: str = "Analyzes a PowerPoint file for basic improvements (slide count, empty placeholders)."

    def _run(self, file_path: str):
        if not os.path.exists(file_path):
            return f"Error: File '{file_path}' not found."
            
        # Check file extension
        if file_path.lower().endswith('.ppt'):
            return (
                "⚠️ This file is in the old PowerPoint format (.ppt). "
                "Please save it as .pptx format using Microsoft PowerPoint or LibreOffice, then upload again. "
                "The python-pptx library only supports the newer .pptx format."
            )
        
        if not file_path.lower().endswith('.pptx'):
            return "Error: Input must be a PowerPoint file (.pptx format)."
        
        try:
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            empty_titles = 0

            for slide in prs.slides:
                # Check if slide has a title and if it's empty
                if slide.shapes.title:
                    if not slide.shapes.title.text.strip():
                        empty_titles += 1
            
            feedback = f"📊 PowerPoint Analysis:\n"
            feedback += f"- Total Slides: {total_slides}\n"
            
            if empty_titles > 0:
                feedback += f"- ⚠️ {empty_titles} slide(s) have empty titles. Consider adding descriptive titles.\n"
            else:
                feedback += "- ✅ All slides have titles.\n"
            
            feedback += f"\n💡 Suggestions: Ensure each slide has clear titles and structured content."
            
            return feedback
        except Exception as e:
            return f"Error analyzing PowerPoint: {str(e)}"

    def _arun(self, file_path: str):
        raise NotImplementedError("Async not implemented")
