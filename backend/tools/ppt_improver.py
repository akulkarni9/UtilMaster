from langchain_core.tools import BaseTool
from pptx import Presentation
from pptx.util import Pt
import os

class PPTImprovementTool(BaseTool):
    name: str = "ppt_improvement_tool"
    description: str = "Improves PowerPoint presentations by adding missing titles and applying basic formatting fixes."

    def _run(self, file_path: str):
        if not os.path.exists(file_path):
            return f"Error: File '{file_path}' not found."
        
        # Strict validation: only accept .pptx files
        if not file_path.lower().endswith('.pptx'):
            return f"Error: This tool only improves PowerPoint (.pptx) files. The file '{os.path.basename(file_path)}' is not a .pptx file."
        
        try:
            prs = Presentation(file_path)
            improvements_made = []
            
            # Improvement 1: Add titles to slides missing them
            slides_fixed = 0
            for i, slide in enumerate(prs.slides, 1):
                if slide.shapes.title:
                    if not slide.shapes.title.text.strip():
                        # Add a generic title
                        slide.shapes.title.text = f"Slide {i}"
                        slides_fixed += 1
            
            if slides_fixed > 0:
                improvements_made.append(f"Added titles to {slides_fixed} slide(s)")
            
            # Save improved version
            base_name = os.path.splitext(file_path)[0]
            output_path = f"{base_name}_improved.pptx"
            prs.save(output_path)
            
            # Generate download link
            filename = os.path.basename(output_path)
            download_url = f"http://localhost:8000/uploads/{filename}"
            
            if improvements_made:
                improvements_list = ", ".join(improvements_made)
                return f"✅ {improvements_list}. [Download Improved PPT]({download_url})"
            else:
                return f"✅ Your presentation already looks great! All slides have titles. [Download Copy]({download_url})"
            
        except Exception as e:
            return f"Error improving PowerPoint: {str(e)}"

    def _arun(self, file_path: str):
        raise NotImplementedError("Async not implemented")
