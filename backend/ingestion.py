import os
from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, UnstructuredPowerPointLoader
from langchain_ollama import OllamaEmbeddings
from langchain_community.vectorstores import Neo4jVector
from langchain_core.documents import Document
import pptx

# Load environment variables
load_dotenv()

class IngestionService:
    def __init__(self):
        self.embeddings = OllamaEmbeddings(
            model="llama3.1:8b",
            base_url=os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
        )
        self.neo4j_url = os.getenv("NEO4J_URI", "bolt://localhost:7687")
        self.neo4j_user = os.getenv("NEO4J_USER", "neo4j")
        self.neo4j_password = os.getenv("NEO4J_PASSWORD", "password")

    def ingest_file(self, file_path: str):
        """Parses a file and stores it in Neo4j."""
        ext = os.path.splitext(file_path)[1].lower()
        docs = []

        try:
            if ext == ".pdf":
                loader = PyPDFLoader(file_path)
                docs = loader.load_and_split()
            elif ext == ".docx":
                loader = Docx2txtLoader(file_path)
                docs = loader.load_and_split()
            elif ext == ".pptx":
                # minimalist pptx loader since Unstructured requires system deps
                # We reuse python-pptx which is already installed
                prs = pptx.Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                
                # Create a single doc for simplicity or chunk it
                full_text = "\n".join(text)
                docs = [Document(page_content=full_text, metadata={"source": file_path})]
            else:
                print(f"Unsupported file type for ingestion: {ext}")
                return

            if docs:
                print(f"Ingesting {len(docs)} chunks from {file_path}...")
                Neo4jVector.from_documents(
                    docs,
                    self.embeddings,
                    url=self.neo4j_url,
                    username=self.neo4j_user,
                    password=self.neo4j_password,
                    index_name="file_embeddings",
                    node_label="Chunk",
                    text_node_property="text",
                    embedding_node_property="embedding"
                )
                print("Ingestion complete.")
                
        except Exception as e:
            print(f"Error ingesting file {file_path}: {e}")
